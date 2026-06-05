from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from app.config import settings
from app.models import AuditProject
from app.services.client_snapshot_service import build_client_snapshot, report_priority_label

SLIDES_DIR = Path(settings.EXPORT_DIR) / "slides"
SLIDES_DIR.mkdir(parents=True, exist_ok=True)


def _safe_filename(value: str) -> str:
    return re.sub(r"[^a-zA-Zа-яА-Я0-9_.-]+", "_", value).strip("_")[:80] or "audit"


def _txt(value: Any) -> str:
    if value is None:
        return "—"
    return str(value)


def _money(value: Any) -> str:
    try:
        return f"{float(value):,.0f} ₽".replace(",", " ")
    except (TypeError, ValueError):
        return "—"


def _add_title(slide, title: str, subtitle: str | None = None):
    slide.shapes.title.text = title
    if subtitle is not None and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _add_bullets(slide, title: str, bullets: list[str]):
    _add_title(slide, title)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(8.4), Inches(4.8))
    frame = box.text_frame
    frame.clear()
    for index, bullet in enumerate(bullets or ["Нет данных"]):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20 if index == 0 else 18)


def generate_pptx_report(project: AuditProject) -> Path:
    """PPTX по client snapshot — те же блоки, что клиентский PDF (R1.8)."""
    snap = build_client_snapshot(project)
    cover = snap.get("cover") or {}
    summary = snap.get("audit_summary") or {}
    metrics = snap.get("metrics") or {}
    offer = snap.get("commercial_offer") or {}

    prs = Presentation()
    prs.core_properties.title = f"PPC-аудит: {cover.get('client_name', 'Клиент')}"
    prs.core_properties.author = "PPC Audit Workspace"

    subtitle = f"{cover.get('niche', '—')} · {snap.get('audit_date', '')}"
    if snap.get("last_analysis_at"):
        subtitle += f"\nАнализ: {snap['last_analysis_at']}"

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(slide, "PPC-аудит и план работ", f"{cover.get('client_name', 'Клиент')}\n{subtitle}")

    if snap.get("is_draft") and snap.get("draft_reason"):
        _add_bullets(prs.slides.add_slide(prs.slide_layouts[1]), "Статус", [snap["draft_reason"]])

    _add_bullets(prs.slides.add_slide(prs.slide_layouts[1]), "Цель аудита", [
        snap.get("goal") or "Не указана",
        f"Сайт: {cover.get('website') or 'не указан'}",
    ])

    _add_bullets(prs.slides.add_slide(prs.slide_layouts[1]), "Краткий вывод", [
        f"Суть: {summary.get('client_problem', '—')}",
        f"Если не исправить: {summary.get('main_risk', '—')}",
        f"Приоритет: {snap.get('priority_label') or report_priority_label(summary.get('priority'))}",
        summary.get("short_conclusion", "—"),
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Ключевые метрики"
    rows = [
        ("Период", _txt(metrics.get("period"))),
        ("Бюджет", _money(metrics.get("budget"))),
        ("Клики", _txt(metrics.get("clicks"))),
        ("Заявки", _txt(metrics.get("leads"))),
        ("Продажи", _txt(metrics.get("sales"))),
        ("Выручка", _money(metrics.get("revenue"))),
        ("CPL", _money(metrics.get("cpl"))),
        ("CPA", _money(metrics.get("cpa"))),
        ("ROMI", f"{metrics.get('romi')}%" if metrics.get("romi") is not None else "—"),
    ]
    table = slide.shapes.add_table(len(rows), 2, Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.8)).table
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(5.4)
    for r, row in enumerate(rows):
        table.cell(r, 0).text = row[0]
        table.cell(r, 1).text = row[1]

    for chart in snap.get("zone_charts") or []:
        data = chart.get("data") or {}
        labels = data.get("labels") or []
        datasets = data.get("datasets") or []
        values = datasets[0].get("data", []) if datasets else []
        bullets = [f"{label}: {value}" for label, value in zip(labels, values)]
        if chart.get("insight"):
            bullets.append(str(chart["insight"]))
        _add_bullets(prs.slides.add_slide(prs.slide_layouts[1]), chart.get("title") or "Приоритет работ", bullets)

    confirmed = snap.get("confirmed_observations") or []
    if confirmed:
        for item in confirmed:
            _add_bullets(prs.slides.add_slide(prs.slide_layouts[1]), f"Наблюдение: {item.get('area') or '—'}", [
                item.get("problem") or "—",
                f"Рекомендация: {item.get('recommendation') or '—'}",
            ])

    services = offer.get("recommended_services") or []
    _add_bullets(prs.slides.add_slide(prs.slide_layouts[1]), "Коммерческое предложение", [
        offer.get("proposal_title", "Аудит и оптимизация Яндекс Директа"),
        f"Срок: {offer.get('estimated_work_days', '—')} дней",
        offer.get("sales_argument", "—"),
        "Услуги: " + "; ".join(services[:5]),
        f"Следующий шаг: {offer.get('next_step', 'согласовать план работ')}",
    ])

    if snap.get("limitations_text"):
        _add_bullets(prs.slides.add_slide(prs.slide_layouts[1]), "Ограничения анализа", [snap["limitations_text"]])

    appendix = snap.get("report_appendix") or []
    if appendix:
        _add_bullets(prs.slides.add_slide(prs.slide_layouts[1]), "Приложение: иллюстрации", [
            "Иллюстрации из рекламных кабинетов с подписью маркетолога.",
        ])
        for item in appendix:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.title.text = item.get("figure_label") or "Иллюстрация"
            image_path = item.get("image_path")
            caption = item.get("caption") or ""
            if image_path and Path(image_path).exists():
                slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.2), width=Inches(8.4))
            box = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(8.4), Inches(1.2))
            box.text_frame.text = caption

    filename = f"ppc_audit_{project.id}_{_safe_filename(cover.get('client_name', 'audit'))}_{datetime.now():%Y%m%d_%H%M%S}.pptx"
    path = SLIDES_DIR / filename
    prs.save(path)
    return path


def upload_pptx_to_google_slides(pptx_path: Path, title: str) -> dict[str, Any]:
    """Опционально загружает PPTX в Google Drive с конвертацией в Google Slides."""
    if not settings.GOOGLE_SERVICE_ACCOUNT_FILE:
        raise RuntimeError(
            "GOOGLE_SERVICE_ACCOUNT_FILE не задан. Скачайте PPTX и импортируйте его в Google Slides вручную "
            "или настройте Google Drive API."
        )

    from google.oauth2 import service_account
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaFileUpload

    scopes = ["https://www.googleapis.com/auth/drive.file"]
    credentials = service_account.Credentials.from_service_account_file(
        settings.GOOGLE_SERVICE_ACCOUNT_FILE, scopes=scopes
    )
    service = build("drive", "v3", credentials=credentials)

    metadata = {
        "name": title,
        "mimeType": "application/vnd.google-apps.presentation",
    }
    if settings.GOOGLE_DRIVE_FOLDER_ID:
        metadata["parents"] = [settings.GOOGLE_DRIVE_FOLDER_ID]

    media = MediaFileUpload(
        str(pptx_path),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=False,
    )
    created = service.files().create(body=metadata, media_body=media, fields="id,name,webViewLink").execute()
    return created
